"""
Generate the RailPulse PowerPoint presentation.
Requires: pip install python-pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Colour palette ──────────────────────────────────────────────
DARK_BG      = RGBColor(0x0B, 0x1D, 0x2F)   # deep navy
MID_BG       = RGBColor(0x12, 0x2E, 0x48)   # medium navy
ACCENT_BLUE  = RGBColor(0x00, 0x9B, 0xDE)   # bright rail blue
ACCENT_TEAL  = RGBColor(0x00, 0xBC, 0xA0)   # teal accent
ACCENT_AMBER = RGBColor(0xF5, 0xA6, 0x23)   # warning / highlight
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY   = RGBColor(0xCC, 0xDC, 0xE8)
DARK_TEXT     = RGBColor(0x1A, 0x2A, 0x3A)
MID_GRAY     = RGBColor(0x7F, 0x95, 0xAF)

# ── Paths ────────────────────────────────────────────────────────
PROJECT_ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
DOCS_DIR     = os.path.join(PROJECT_ROOT, "docs")
ASSETS_DIR   = os.path.join(DOCS_DIR, "assets")
PROJ_ASSETS   = os.path.join(PROJECT_ROOT, "project-instructions", "assets")
OUTPUT_PATH  = os.path.join(DOCS_DIR, "RailPulse_Presentation.pptx")

IMAGE_MAP = {
    "erd":        os.path.join(ASSETS_DIR, "erd-readme.png"),
    "dashboard_q1": os.path.join(ASSETS_DIR, "dashboard-q1-peak-hour.png"),
    "dashboard_q2": os.path.join(ASSETS_DIR, "dashboard-q2-platforms.png"),
    "dashboard_q3": os.path.join(ASSETS_DIR, "dashboard-q3-destinations.png"),
    "dashboard_q4": os.path.join(ASSETS_DIR, "dashboard-q4-frequency.png"),
    "dashboard_q5": os.path.join(ASSETS_DIR, "dashboard-q5-accessibility.png"),
    "dashboard_leaderboard": os.path.join(ASSETS_DIR, "dashboard-leaderboard.png"),
    "dashboard_dataquality": os.path.join(ASSETS_DIR, "dashboard-data-quality.png"),
    "dashboard_overview": os.path.join(ASSETS_DIR, "dashboard-overview.png"),
    "sql_chat":     os.path.join(ASSETS_DIR, "dashboard-sql-chat.png"),
    "architecture": os.path.join(PROJ_ASSETS, "project_architecture.png"),
    "dev_portal":   os.path.join(PROJ_ASSETS, "developer_portal.png"),
}

# ── Helpers ──────────────────────────────────────────────────────
def _slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def _add_textbox(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Calibri", line_spacing=1.15):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(4)
    if line_spacing:
        p.line_spacing = Pt(font_size * line_spacing)
    return tf

def _add_rich_box(slide, left, top, width, height, paragraphs,
                  default_size=14, default_color=WHITE, default_name="Calibri"):
    """paragraphs: list of (text, font_size, color, bold, alignment) tuples."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, pdata in enumerate(paragraphs):
        text = pdata[0]
        size = pdata[1] if len(pdata) > 1 and pdata[1] else default_size
        color = pdata[2] if len(pdata) > 2 and pdata[2] else default_color
        bold = pdata[3] if len(pdata) > 3 and pdata[3] else False
        align = pdata[4] if len(pdata) > 4 and pdata[4] else PP_ALIGN.LEFT
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = default_name
        p.alignment = align
        p.space_after = Pt(4)
    return tf

def _add_title_bar(slide, title_text, subtitle_text=None):
    """Dark blue title strip at the top."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        Inches(13.333), Inches(1.15)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = MID_BG
    shape.line.fill.background()

    _add_textbox(slide, 0.7, 0.15, 11.5, 0.55, title_text,
                 font_size=30, color=WHITE, bold=True)
    if subtitle_text:
        _add_textbox(slide, 0.7, 0.7, 11.5, 0.35, subtitle_text,
                     font_size=14, color=MID_GRAY, bold=False)

def _add_page_number(slide, num):
    _add_textbox(slide, 12.2, 7.1, 0.9, 0.3, str(num),
                 font_size=10, color=MID_GRAY, alignment=PP_ALIGN.RIGHT)

def _add_accent_line(slide, left, top, width):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(left), Inches(top),
        Inches(width), Inches(0.04)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_BLUE
    shape.line.fill.background()

def _add_bullet_card(slide, left, top, width, height, title, items, title_color=ACCENT_BLUE):
    """A card with a coloured title and bulleted items."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top),
        Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = MID_BG
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(14)
    p.font.color.rgb = title_color
    p.font.bold = True
    p.font.name = "Calibri"
    p.space_after = Pt(6)
    for item in items:
        p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(11)
        p.font.color.rgb = LIGHT_GRAY
        p.font.name = "Calibri"
        p.space_after = Pt(2)
        p.level = 1

def _add_image_safe(slide, key, left, top, width, height=None):
    path = IMAGE_MAP.get(key)
    if path and os.path.exists(path):
        if height:
            return slide.shapes.add_picture(path, Inches(left), Inches(top),
                                           Inches(width), Inches(height))
        else:
            return slide.shapes.add_picture(path, Inches(left), Inches(top),
                                           Inches(width))
    return None

def _add_section_number_shape(slide, left, top, num):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(left), Inches(top), Inches(0.55), Inches(0.55)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_BLUE
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = str(num)
    p.font.size = Pt(18)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER

def _add_kpi_tile(slide, left, top, value, label, color=ACCENT_BLUE):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top),
        Inches(2.0), Inches(1.2)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = MID_BG
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = str(value)
    p.font.size = Pt(32)
    p.font.color.rgb = color
    p.font.bold = True
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(2)
    p2 = tf.add_paragraph()
    p2.text = label
    p2.font.size = Pt(10)
    p2.font.color.rgb = LIGHT_GRAY
    p2.font.name = "Calibri"
    p2.alignment = PP_ALIGN.CENTER

# ══════════════════════════════════════════════════════════════════
#  BUILD PRESENTATION
# ══════════════════════════════════════════════════════════════════
prs = Presentation()
prs.slide_width  = Inches(13.333)   # 16:9 widescreen
prs.slide_height = Inches(7.5)

SW = prs.slide_width   # 13.333
SH = prs.slide_height  # 7.5

# ── Helper: layout constants ──
CONTENT_TOP = 1.45
CONTENT_LM = 0.7

# ═══════════════════════  SLIDE 1 – TITLE  ══════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
_slide_bg(slide, DARK_BG)

# large accent shape top-right
shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(8.5), Inches(0), Inches(4.833), Inches(7.5)
)
shape.fill.solid()
shape.fill.fore_color.rgb = MID_BG
shape.line.fill.background()

_add_image_safe(slide, "architecture", 0.5, 1.5, 7.5)

# Title
_add_textbox(slide, 0.7, 0.6, 12.0, 1.0, "RailPulse",
             font_size=48, color=WHITE, bold=True)
_add_accent_line(slide, 0.7, 1.5, 2.5)
_add_textbox(slide, 0.7, 1.7, 9.0, 1.0,
             "Belgian Transit SQL Analysis",
             font_size=28, color=ACCENT_BLUE, bold=False)

_add_textbox(slide, 9.2, 1.5, 3.8, 5.5,
             "Sprint 1 — Local Ingestion & Relational Modelling\n\n"
             "Enterprise-grade ELT pipeline over the SNCB/NMBS national timetable\n\n"
             "2.17 million scheduled departures\n"
             "652 stations · 1,801 routes\n"
             "9 data-quality rules · 47 labelled queries\n"
             "Streamlit dashboard + text-to-SQL preview",
             font_size=14, color=LIGHT_GRAY)

_add_textbox(slide, 0.7, 5.5, 8.0, 0.8,
             "Stéphane van der Aa  ·  BeCode AI & Data Science  ·  24 July 2026",
             font_size=13, color=MID_GRAY)
_add_page_number(slide, 1)

# ═══════════════════════  SLIDE 2 – AGENDA  ═════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_slide_bg(slide, DARK_BG)
_add_title_bar(slide, "Agenda")
_add_section_number_shape(slide, 0.7, 1.55, "📋")

agenda_items = [
    ("01", "Mission & Context", "RailPulse consulting firm, SNCB client, winter scheduling optimisation"),
    ("02", "Data & Architecture", "GTFS feed, ELT pipeline, SQLite 3NF model, repository structure"),
    ("03", "Data Quality", "9 DQ rules, quarantine system, 3 feed quirks that shaped the analysis"),
    ("04", "The Five Analytical Questions", "Peak hour, platform bottlenecks, morning destinations, service frequency, accessibility"),
    ("05", "Key Findings", "Bruxelles-Central bottleneck, annualisation insight, accessibility gap"),
    ("06", "Nice-to-Haves & Dashboard", "Network leaderboard, index optimisation, Streamlit dashboard, SQL Chat"),
    ("07", "Database Theory", "Key interview topics from the study guide"),
    ("08", "Project Deliverables & Timeline", "Evaluation criteria, timeline, deliverables checklist"),
]

for i, (num, title, desc) in enumerate(agenda_items):
    row = i // 2
    col = i % 2
    left = 0.7 + col * 6.1
    top = 2.1 + row * 1.25
    _add_textbox(slide, left, top, 0.5, 0.35, num, font_size=20, color=ACCENT_BLUE, bold=True)
    _add_textbox(slide, left + 0.6, top, 5.2, 0.35, title, font_size=16, color=WHITE, bold=True)
    _add_textbox(slide, left + 0.6, top + 0.33, 5.2, 0.65, desc, font_size=11, color=MID_GRAY)

_add_accent_line(slide, 0.7, 7.15, 2.0)
_add_page_number(slide, 2)

# ═══════════════════════  SLIDE 3 – MISSION & CONTEXT  ══════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_slide_bg(slide, DARK_BG)
_add_title_bar(slide, "Mission & Context", "Section 01")
_add_section_number_shape(slide, 0.7, 1.55, "01")

_add_textbox(slide, CONTENT_LM, 1.7, 6.0, 0.4, "The Mission",
             font_size=20, color=ACCENT_BLUE, bold=True)
_add_textbox(slide, CONTENT_LM, 2.15, 10.0, 1.8,
             "RailPulse is an urban mobility consulting firm. The Belgian National "
             "Railway company (SNCB/NMBS) wants a clear overview of operational "
             "performance and delay patterns to optimise their winter scheduling. "
             "Our mission is to extract liveboard data directly from the iRail API, "
             "build a normalized database, and provide an analytical report detailing "
             "network bottlenecks.",
             font_size=14, color=LIGHT_GRAY)

_add_bullet_card(slide, 0.7, 4.2, 5.8, 2.7, "Learning Objectives", [
    "Design and build an SQL database schema from API data structures",
    "Query relational tables to extract performance metrics",
    "Handle date/time objects and calculate delays in SQL",
    "Present an operational efficiency analysis to a business client",
    "Consolidate: JOINs, GROUP BY, aggregations, time-series in SQLite",
])

_add_bullet_card(slide, 6.8, 4.2, 5.8, 2.7, "Challenge Parameters", [
    "Duration: 4 days · Deadline: 24/07/2026 5:00 PM",
    "Team challenge (in spirit), individual delivery",
    "NO pandas for filtering/aggregating — Python only for network I/O + sqlite3",
    "All analytical questions solved with standard SQL operations only",
    "Team feedback: 5 min overview + random Q&A from study guide (SQL&DB_theory.md)",
], title_color=ACCENT_AMBER)

_add_page_number(slide, 3)

# ═══════════════════════  SLIDE 4 – DATA & THE FEED  ═════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_slide_bg(slide, DARK_BG)
_add_title_bar(slide, "Data & The GTFS Feed", "Section 02")
_add_section_number_shape(slide, 0.7, 1.55, "02")

_add_textbox(slide, CONTENT_LM, 1.6, 11.8, 0.4, "Data Source: Belgian Mobility Open Data Portal",
             font_size=18, color=ACCENT_BLUE, bold=True)
_add_textbox(slide, CONTENT_LM, 2.05, 11.8, 1.4,
             "We ingest GTFS Static (26 MB ZIP) and GTFS Realtime (JSON, every 30 s) from the "
             "official SNCB/NMBS developer portal. The static feed covers the full national timetable "
             "for 358 operating dates (2025-12-20 → 2026-12-12). A free Standard subscription key "
             "provides 12,000 requests/day at 500/min — our pipeline uses 1 request per rebuild.",
             font_size=13, color=LIGHT_GRAY)

_add_image_safe(slide, "dev_portal", 0.7, 3.6, 5.5, 3.4)

# KPI tiles on the right
_add_kpi_tile(slide, 7.0, 3.6, "2.17M", "Timetabled Calls")
_add_kpi_tile(slide, 9.3, 3.6, "134,809", "Trips")
_add_kpi_tile(slide, 11.6, 3.6, "652", "Stations")
_add_kpi_tile(slide, 7.0, 5.1, "1,801", "Routes")
_add_kpi_tile(slide, 9.3, 5.1, "4.7M", "Service Dates")
_add_kpi_tile(slide, 11.6, 5.1, "1 GB", "Database")

_add_page_number(slide, 4)

# ═══════════════════════  SLIDE 5 – PROJECT ROADMAP  ═════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_slide_bg(slide, DARK_BG)
_add_title_bar(slide, "Project Roadmap — Four Sprints", "Section 02")
_add_section_number_shape(slide, 0.7, 1.55, "02")

sprints = [
    ("Sprint 1\nLocal Ingestion &\nRelational Modelling", ACCENT_BLUE,
     ["SQLite database + 3NF schema", "GTFS Static feed ingestion", "9 DQ rules, quarantine system",
      "5 analytical questions (SQL only)", "Streamlit dashboard, 47 queries"]),
    ("Sprint 2\nCloud Migration & Serverless\nPipelines", RGBColor(0xE7, 0x4C, 0x3C),
     ["Azure SQL / PostgreSQL", "Azure Functions, Timer Trigger CRON",
      "Poll GTFS-RT trip updates + alerts", "Every 15-30 min, budget-controlled",
      "Environment variables for secure connections"]),
    ("Sprint 3\nEnterprise BI\n(Power BI)", RGBColor(0xF3, 0x9C, 0x12),
     ["Power BI Desktop connected to Azure", "Semantic data model (star/snowflake)",
      "DAX measures: On-Time Rate %, peak matrices", "Interactive drill-down, bookmarks",
      "Cross-hub comparison dashboards"]),
    ("Sprint 4\nConversational Transit\nAssistant (GenAI)", RGBColor(0x27, 0xAE, 0x60),
     ["Local LLM via Ollama / HuggingFace", "Text-to-SQL with schema-aware prompts",
      "Regex guardrails: block DROP, DELETE", "Execution caps: timeout + row limit",
      "Streamlit chat interface — preview included!"]),
]

cards_per_row = 4
card_w = 2.95
card_h = 4.4
gap = 0.25
start_left = 0.55

for i, (title, color, items) in enumerate(sprints):
    left = start_left + i * (card_w + gap)
    top = 1.6
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top),
        Inches(card_w), Inches(card_h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = MID_BG
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True

    # Sprint number bar
    p = tf.paragraphs[0]
    p.text = f"SPRINT {i+1}"
    p.font.size = Pt(9)
    p.font.color.rgb = color
    p.font.bold = True
    p.font.name = "Calibri"
    p.space_after = Pt(2)

    p = tf.add_paragraph()
    p.text = title
    p.font.size = Pt(13)
    p.font.color.rgb = color
    p.font.bold = True
    p.font.name = "Calibri"
    p.space_after = Pt(8)

    # accent line
    p = tf.add_paragraph()
    p.text = "─" * 24
    p.font.size = Pt(6)
    p.font.color.rgb = color
    p.space_after = Pt(6)

    for item in items:
        p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(9)
        p.font.color.rgb = LIGHT_GRAY
        p.font.name = "Calibri"
        p.space_after = Pt(3)

_add_textbox(slide, 0.7, 6.3, 11.5, 0.8,
             "🛠️ Tools: Python, SQLite / PostgreSQL, Azure, Power BI, Ollama, LangChain, Streamlit",
             font_size=12, color=MID_GRAY)

_add_page_number(slide, 5)

# ═══════════════════════  SLIDE 6 – ARCHITECTURE  ════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_slide_bg(slide, DARK_BG)
_add_title_bar(slide, "Pipeline Architecture — ELT, not ETL", "Section 02")
_add_section_number_shape(slide, 0.7, 1.55, "02")

_add_textbox(slide, CONTENT_LM, 1.6, 11.5, 1.0,
             "Python does network I/O and executes SQL. Every cleaning rule, every constraint, and "
             "every metric lives in a .sql file. Python never inspects a value — raw CSV rows land "
             "verbatim in staging tables (all TEXT, no constraints), and the transform layer decides "
             "what is clean, what is quarantined, and why.",
             font_size=13, color=LIGHT_GRAY)

_add_image_safe(slide, "architecture", 0.7, 2.8, 7.0, 4.0)

_add_bullet_card(slide, 8.2, 2.8, 4.5, 4.3, "Pipeline Stages", [
    "stg_* (verbatim CSV → TEXT columns, no constraints)",
    "02_schema.sql (3NF core: 11 entity + 6 ref_ tables, 20 FKs)",
    "03_transform.sql (9 DQ rules, quarantine → rejected_row)",
    "04_indexes.sql (1 index per named query, built post-load)",
    "05_views.sql (semantic layer: v_departure, v_service_frequency, …)",
    "06_realtime.sql (GTFS-RT additive tables, survives rebuild)",
    "07_cleanup.sql (drop staging, VACUUM — database shrinks to ~1 GB)",
    "sql/analysis/q1..q7.sql (47 labelled queries → output/*.csv)",
])

_add_page_number(slide, 6)

# ═══════════════════════  SLIDE 7 – REPOSITORY STRUCTURE  ═══════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_slide_bg(slide, DARK_BG)
_add_title_bar(slide, "Repository Structure", "Section 02")
_add_section_number_shape(slide, 0.7, 1.55, "02")

_add_textbox(slide, CONTENT_LM, 1.6, 11.8, 0.8,
             "Every rule and every metric lives under sql/. Python has exactly two jobs: network I/O and "
             "executing SQL. Nothing in src/ inspects, filters or aggregates a data value.",
             font_size=13, color=LIGHT_GRAY)

structure_items = [
    ("sql/", "schema, transform, indexes, views, analysis — every number traces to a .sql file"),
    ("src/railpulse/", "thin Python layer: build, analyse, verify, benchmark, ingest, poll"),
    ("scripts/", "poll_realtime.sh (cron/launchd), setup_api_key.py (Playwright)"),
    ("tests/", "149 tests over a synthetic broken feed — each DQ rule tested"),
    ("dashboard/", "Streamlit report + optional SQL Chat (text-to-SQL, local model)"),
    ("docs/", "ERD, data dictionary, analysis report, ADRs, data quality, glossary, API compliance"),
    ("output/", "generated: one CSV per analysis query"),
]

for i, (label, desc) in enumerate(structure_items):
    top = 2.6 + i * 0.65
    _add_textbox(slide, 0.7, top, 2.8, 0.45, label, font_size=14, color=ACCENT_BLUE, bold=True)
    _add_textbox(slide, 3.6, top + 0.02, 9.0, 0.45, desc, font_size=12, color=LIGHT_GRAY)

_add_page_number(slide, 7)

# ═══════════════════════  SLIDE 8 – ERD  ═════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_slide_bg(slide, DARK_BG)
_add_title_bar(slide, "Entity Relationship Diagram", "Section 02")
_add_section_number_shape(slide, 0.7, 1.55, "02")

_add_textbox(slide, CONTENT_LM, 1.6, 11.5, 0.8,
             "3NF model with 11 entity tables, 6 ref_ lookup tables and 20 enforced foreign keys. "
             "station and platform are split (not a single stops table). stop_time is the fact table at "
             "2,165,507 rows — event grain: one scheduled call of one trip at one platform.",
             font_size=13, color=LIGHT_GRAY)

_add_image_safe(slide, "erd", 0.7, 2.6, 7.5, 4.5)

_add_bullet_card(slide, 8.6, 2.6, 4.2, 4.5, "Key Modelling Decisions", [
    "station / platform split (2 grains in one GTFS file → 2 tables)",
    "service keeps weekday columns + has_weekday_pattern flag (DQ-01)",
    "stop_time materialises departure_hour, is_boardable (SARGability)",
    "WITHOUT ROWID on service_date saves ~190 MB",
    "Every GTFS code column → FOREIGN KEY into ref_ table",
    "rt_trip_update.trip_id is a soft link (not FK) by design",
    "PRAGMA foreign_key_check returns clean",
])

_add_page_number(slide, 8)

# ═══════════════════════  SLIDE 9 – DATA QUALITY  ════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_slide_bg(slide, DARK_BG)
_add_title_bar(slide, "Data Quality — 9 Rules, 12 Quarantined Rows", "Section 03")
_add_section_number_shape(slide, 0.7, 1.55, "03")

_add_textbox(slide, CONTENT_LM, 1.65, 11.8, 1.0,
             "7,057,090 rows staged → 7,057,078 loaded → 12 quarantined (0.00017%). "
             "Nothing is silently dropped. Each rejected row lands in rejected_row with its rule, "
             "reason, source file, physical line number and a JSON snapshot of the payload.",
             font_size=13, color=LIGHT_GRAY)

dq_rules = [
    ("DQ-01", "calendar.txt empty", "All 51,593 services have zeroed weekday flags — the real calendar is 4.7M rows in calendar_dates.txt"),
    ("DQ-02", "Code 0 ≠ 'no'", "wheelchair/bike code 0 means 'no information', not refusal. ref_accessibility.is_guaranteed encodes the distinction"),
    ("DQ-03", "Implausible times", "12 calls at 63:18–87:39 into service day → quarantined. 2 trips affected, truncated in core model"),
    ("DQ-04", "Orphan FKs", "0 rows — screen before FK enforcement to convert fatal error → quarantined row"),
    ("DQ-05", "Duplicate keys", "0 rows — ROW_NUMBER() partition guard + audit pass on INSERT OR IGNORE"),
]

for i, (code, title, desc) in enumerate(dq_rules):
    left = 0.45 + (i % 3) * 4.2
    top = 2.85 + (i // 3) * 2.0
    _add_textbox(slide, left, top, 3.8, 0.3, f"{code} — {title}", font_size=12, color=ACCENT_BLUE, bold=True)
    _add_textbox(slide, left, top + 0.32, 3.8, 1.5, desc, font_size=10, color=LIGHT_GRAY)

_add_textbox(slide, 0.5, 6.8, 12.0, 0.4,
             "Three feed quirks that changed the answers: (1) empty calendar.txt weekday flags, "
             "(2) 577K pass-through calls, (3) 31K calls at ≥ 24:00:00  —  each requires explicit handling",
             font_size=11, color=ACCENT_AMBER)

_add_page_number(slide, 9)

# ═══════════════════════  SLIDE 10 – THE FIVE QUESTIONS OVERVIEW ══
slide = prs.slides.add_slide(prs.slide_layouts[6])
_slide_bg(slide, DARK_BG)
_add_title_bar(slide, "The Five Analytical Questions", "Section 04")
_add_section_number_shape(slide, 0.7, 1.55, "04")

questions = [
    ("Q1", "The Peak Hour Problem", "What hour experiences the highest volume of scheduled departures?",
     "17:00–17:59  ·  950,651 annual departures"),
    ("Q2", "Platform Bottlenecks", "Top 3 busiest platforms at Bruxelles-Central?",
     "Platforms 4, 3, 2  ·  up to 12.8 trains/hour"),
    ("Q3", "Morning Destinations", "Top 3 terminal destinations for trips before 12:00?",
     "Anvers-Central, Louvain, Charleroi-Central"),
    ("Q4", "Service Frequency", "Classify each service: High (5+), Medium (2-4), Low (1) d/wk?",
     "45.2% High · 34.7% Medium · 20.1% Low"),
    ("Q5", "Accessibility Audit", "Ratio of trips guaranteeing wheelchair/bike?",
     "91.3% bike · wheelchair field empty for all 134,809 trips"),
]

for i, (num, title, q, a) in enumerate(questions):
    top = 1.7 + i * 1.1
    _add_textbox(slide, 0.7, top, 0.5, 0.35, num, font_size=22, color=ACCENT_BLUE, bold=True)
    _add_textbox(slide, 1.25, top, 5.0, 0.35, title, font_size=16, color=WHITE, bold=True)
    _add_textbox(slide, 1.25, top + 0.33, 6.5, 0.35, q, font_size=11, color=MID_GRAY)
    _add_textbox(slide, 8.5, top + 0.1, 4.5, 0.6, f"↳ {a}", font_size=12, color=ACCENT_AMBER)

_add_textbox(slide, CONTENT_LM, 7.0, 11.5, 0.3,
             "All answers solved with SQL only — no Python filtering or aggregation.",
             font_size=11, color=MID_GRAY, alignment=PP_ALIGN.CENTER)
_add_page_number(slide, 10)

# ═══════════════════════  SLIDE 11 – Q1 OVERVIEW  ════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_slide_bg(slide, DARK_BG)
_add_title_bar(slide, "Q1 — The Peak Hour Problem", "Section 04 · Question 1")
_add_section_number_shape(slide, 0.7, 1.55, "Q1")

_add_textbox(slide, 1.35, 1.6, 10.0, 0.4,
             "Answer: 17:00–17:59 — 950,651 annual departures (6.51% of network day)",
             font_size=18, color=ACCENT_AMBER, bold=True)

_add_textbox(slide, CONTENT_LM, 2.2, 12.0, 1.8,
             "The obvious query — SELECT departure_hour, COUNT(*) ... GROUP BY departure_hour — gives a "
             "confidently wrong answer: 10:00. The problem is that counting rows counts timetable variants, "
             "not trains. A trip that runs once and a trip that runs 250 times are one row each. "
             "Weighting each call by v_trip_service_days.operating_days (annualisation) moves hour 17 "
             "from rank 10 to rank 1. The gap between the two rankings IS the finding.",
             font_size=13, color=LIGHT_GRAY)

_add_image_safe(slide, "dashboard_q1", 0.7, 4.1, 7.0, 3.0)

_add_bullet_card(slide, 8.1, 4.1, 4.7, 3.0, "How Ranks Shift", [
    "Naïve rank 1: 10:00 (94,323 timetable rows)",
    "Annualised rank 1: 17:00 (950,651 departures)",
    "Hour 17: rank 10 → rank 1  (+9)",
    "Hour 07: rank 12 → rank 2 (+10)",
    "Hour 10: rank 1  → rank 8 (−7)",
    "Midday calls avg 8.8 days · Evening avg 12.4 days",
    "Seasonal services over-represented in raw count",
])

_add_page_number(slide, 11)

# ═══════════════════════  SLIDE 12 – Q2 OVERVIEW  ════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_slide_bg(slide, DARK_BG)
_add_title_bar(slide, "Q2 — Platform Bottlenecks at Bruxelles-Central", "Section 04 · Question 2")
_add_section_number_shape(slide, 0.7, 1.55, "Q2")

_add_textbox(slide, 1.35, 1.6, 10.0, 0.4,
             "Answer: Platforms 4 (63,426), 3 (62,276), 2 (56,874) annual departures",
             font_size=18, color=ACCENT_AMBER, bold=True)

_add_textbox(slide, CONTENT_LM, 2.15, 11.5, 1.0,
             "The top three platforms are robust across both measures (annualised and raw), but the order "
             "between platforms 3 and 4 swaps — platform 3 leads on raw timetable rows, platform 4 leads "
             "when weighted by operating days. 1,348 additional calls (2.69%) carry no platform allocation.",
             font_size=13, color=LIGHT_GRAY)

_add_image_safe(slide, "dashboard_q2", 0.7, 3.3, 7.0, 3.8)

# Right side table
_add_textbox(slide, 8.2, 3.3, 4.5, 0.3, "Platform peak-hour pressure",
             font_size=14, color=ACCENT_BLUE, bold=True)
peak_data = "Platform 3 · 08:00 · 12.8 trains/day · one every 4.7 min\n" \
            "Platform 4 · 16:00 · 11.3 trains/day · one every 5.3 min\n" \
            "Platform 2 · 07:00 · 10.9 trains/day · one every 5.5 min\n" \
            "Platform 6 · 16:00 ·  8.5 trains/day · one every 7.1 min"
_add_textbox(slide, 8.2, 3.7, 4.5, 3.0, peak_data, font_size=11, color=LIGHT_GRAY)

_add_page_number(slide, 12)

# ═══════════════════════  SLIDE 13 – Q3 OVERVIEW  ════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_slide_bg(slide, DARK_BG)
_add_title_bar(slide, "Q3 — Busiest Morning Destinations", "Section 04 · Question 3")
_add_section_number_shape(slide, 0.7, 1.55, "Q3")

_add_textbox(slide, 1.35, 1.6, 10.0, 0.4,
             "Answer: Anvers-Central (41,972), Louvain (27,516), Charleroi-Central (21,328)",
             font_size=18, color=ACCENT_AMBER, bold=True)

_add_textbox(slide, CONTENT_LM, 2.15, 11.5, 1.0,
             "'Morning' = trip that originates before 12:00:00. The ranking uses annualised trips "
             "(trips × operating days), not distinct trip IDs. Bruxelles-Midi is 2nd by distinct "
             "services but 4th by annualised trips — its morning services average only 5.3 operating days.",
             font_size=13, color=LIGHT_GRAY)

_add_image_safe(slide, "dashboard_q3", 0.7, 3.3, 7.0, 3.8)

_add_textbox(slide, 8.2, 3.3, 4.5, 0.3, "Most morning-skewed destinations",
             font_size=14, color=ACCENT_BLUE, bold=True)
skew_data = "1. Schaerbeek — 65.8% of arrivals before noon\n" \
           "2. Eupen — 54.7%\n" \
           "3. Blankenberge — 51.7%\n" \
           "4. Wavre — 49.9%\n" \
           "5. Lille Flandres (FR) — 49.2%\n\n" \
           "Schaerbeek is a depot + maintenance hub — consistent with strong morning positioning flow."
_add_textbox(slide, 8.2, 3.7, 4.5, 3.0, skew_data, font_size=11, color=LIGHT_GRAY)

_add_page_number(slide, 13)

# ═══════════════════════  SLIDE 14 – Q4 OVERVIEW  ════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_slide_bg(slide, DARK_BG)
_add_title_bar(slide, "Q4 — Service Frequency Classification", "Section 04 · Question 4")
_add_section_number_shape(slide, 0.7, 1.55, "Q4")

_add_textbox(slide, 1.35, 1.6, 10.0, 0.4,
             "Answer: 45.2% High · 34.7% Medium · 20.1% Low/Special  (but the 45% carries 86.2% of operating days)",
             font_size=18, color=ACCENT_AMBER, bold=True)

_add_textbox(slide, CONTENT_LM, 2.15, 11.5, 1.2,
             "A data problem had to be solved first: calendar.txt weekday flags are ALL ZERO for all 51,593 "
             "services (DQ-01). Without that fix, the textbook CASE WHEN query classifies 100% of the Belgian "
             "rail network as 'Low Frequency'. The weekly rhythm was derived from the 4.7M rows of "
             "calendar_dates.txt via v_service_frequency (modal days per active week).",
             font_size=13, color=LIGHT_GRAY)

_add_image_safe(slide, "dashboard_q4", 0.7, 3.5, 7.0, 3.5)

_add_textbox(slide, 8.2, 3.5, 4.5, 0.3, "Definition sensitivity",
             font_size=14, color=ACCENT_BLUE, bold=True)
sens_data = "'High Frequency' share by definition:\n\n" \
           "A. Modal days per active week → 45.24%\n" \
           "B. Distinct weekdays ever touched → 62.26%\n" \
           "C. Busiest single week → 52.10%\n\n" \
           "Headline moves by 17 pp depending on definition.\n" \
           "All three published — single number would\nimply precision the question doesn't have."
_add_textbox(slide, 8.2, 3.9, 4.5, 3.0, sens_data, font_size=11, color=LIGHT_GRAY)

_add_page_number(slide, 14)

# ═══════════════════════  SLIDE 15 – Q5 OVERVIEW  ════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_slide_bg(slide, DARK_BG)
_add_title_bar(slide, "Q5 — The Accessibility Audit", "Section 04 · Question 5")
_add_section_number_shape(slide, 0.7, 1.55, "Q5")

_add_textbox(slide, 1.35, 1.6, 10.0, 0.4,
             "Bicycle storage: 91.3% guaranteed · Wheelchair: unpopulated for all 134,809 trips",
             font_size=18, color=ACCENT_AMBER, bold=True)

_add_textbox(slide, CONTENT_LM, 2.15, 11.5, 1.2,
             "The bicycle gap is a mode gap, not a route gap: 100% of 123,051 rail trips guarantee bike "
             "storage, 0% of 11,758 rail-replacement bus trips do. All 270 zero-scoring routes are buses. "
             "wheelchair_accessible is empty for ALL 134,809 trips — not a '0% accessible' finding, but a "
             "publishing gap: no journey planner using this feed can answer a wheelchair question at all.",
             font_size=13, color=LIGHT_GRAY)

_add_image_safe(slide, "dashboard_q5", 0.7, 3.5, 7.0, 3.5)

_add_textbox(slide, 8.2, 3.5, 4.5, 0.3, "Recommendations",
             font_size=14, color=ACCENT_BLUE, bold=True)
rec_data = "1. Populate wheelchair_accessible at trip level\n" \
          "   and wheelchair_boarding at station level\n" \
          "   (single-digit cardinality fields, cheap to fix)\n\n" \
          "2. Set bikes_allowed on replacement buses\n" \
          "   explicitly (0=unknown, 2=no) — currently\n" \
          "   indistinguishable from an unset field\n\n" \
          "3. Procurement decision: bicycle capacity in\n" \
          "   the replacement-bus contract — not a\n" \
          "   route-by-route investigation"
_add_textbox(slide, 8.2, 3.9, 4.5, 3.0, rec_data, font_size=10, color=LIGHT_GRAY)

_add_page_number(slide, 15)

# ═══════════════════════  SLIDE 16 – KEY FINDINGS  ═══════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_slide_bg(slide, DARK_BG)
_add_title_bar(slide, "Key Findings — The Three Things Worth Acting On", "Section 05")
_add_section_number_shape(slide, 0.7, 1.55, "05")

findings = [
    ("1", "Bruxelles-Central is the network's structural bottleneck",
     ACCENT_AMBER,
     "Handles more annual departures than Bruxelles-Midi (311,324 vs 283,415) across "
     "6 platforms instead of 21 — that is 8,113 timetabled calls per platform against "
     "Midi's 2,085, a 3.9× pressure differential. Its busiest platform turns over "
     "12.8 trains in its peak hour, one every 4.7 minutes. No trough in the day to absorb disruption."),

    ("2", "The evening peak is real and the timetable file hides it",
     ACCENT_BLUE,
     "Naïve COUNT(*) says the network peaks at 10:00. Annualised departures say 17:00. "
     "Hour 17 sits at rank 10 on the naïve measure and rank 1 on the real one. "
     "Capacity decisions taken from an unweighted timetable count would invest in the wrong hour. "
     "The mechanism: SNCB fragments each train number into ~21 trip rows, each covering ~9 dates."),

    ("3", "The accessibility data cannot support an accessibility statement",
     ACCENT_TEAL,
     "Not because the network performs badly, but because the field is empty. "
     "wheelchair_accessible carries no value for any of 134,809 trips. "
     "wheelchair_boarding carries no value for any of 652 stations. "
     "This is a publishing gap — cheap to fix relative to what it blocks."),
]

for i, (num, title, color, desc) in enumerate(findings):
    top = 1.6 + i * 1.85
    _add_textbox(slide, 0.7, top, 1.5, 0.4, f"Finding {num}", font_size=16, color=color, bold=True)
    _add_textbox(slide, 0.7, top + 0.35, 11.8, 0.35, title, font_size=18, color=WHITE, bold=True)
    _add_textbox(slide, 0.7, top + 0.75, 11.8, 0.9, desc, font_size=12, color=LIGHT_GRAY)

_add_page_number(slide, 16)

# ═══════════════════════  SLIDE 17 – NICE-TO-HAVES  ═════════════=
slide = prs.slides.add_slide(prs.slide_layouts[6])
_slide_bg(slide, DARK_BG)
_add_title_bar(slide, "Nice-to-Have Features", "Section 06")
_add_section_number_shape(slide, 0.7, 1.55, "06")

# Live Stream
_add_textbox(slide, CONTENT_LM, 1.7, 5.0, 0.35, "Live Stream Integration",
             font_size=18, color=ACCENT_BLUE, bold=True)
_add_textbox(slide, CONTENT_LM, 2.1, 5.0, 1.2,
             "GTFS-RT trip updates and service alerts polled every 30 s via scripts/poll_realtime.sh "
             "(cron / launchd). Append-only rt_* tables survive rebuilds. Idempotent: UNIQUE(feed, "
             "feed_timestamp_epoch) so over-polling skips duplicates. 100% of real-time trip IDs "
             "resolved against the static feed in observed window.",
             font_size=12, color=LIGHT_GRAY)

# Network Leaderboard
_add_textbox(slide, CONTENT_LM, 3.5, 5.0, 0.35, "Network Leaderboard",
             font_size=18, color=ACCENT_BLUE, bold=True)
_add_textbox(slide, CONTENT_LM, 3.9, 5.0, 1.0,
             "Five main hubs compared: Bruxelles-Central, -Nord, -Midi, Anvers-Central, Gand-Saint-Pierre. "
             "Composite score (connectivity, platform headroom, load smoothness) + real-time punctuality. "
             "Punctuality from live GTFS-RT: observed delay distribution, on-time (<2 min) rates, "
             "cancellation tracking — mechanism works, needs a week of data for stable ranking.",
             font_size=12, color=LIGHT_GRAY)

# Index optimization
_add_textbox(slide, CONTENT_LM, 5.1, 5.0, 0.35, "Index Optimisation",
             font_size=18, color=ACCENT_BLUE, bold=True)
_add_textbox(slide, CONTENT_LM, 5.45, 5.0, 1.2,
             "EXPLAIN QUERY PLAN on every analytical query. Measured benchmarks: SARGable violations cost "
             "~100× (Q1 histogram) to ~500× (Q2 platform lookup). Each index justified by a named query. "
             "Composite indexes led by the equality predicate, grouping column following. "
             "make benchmark reproduces all measurements.",
             font_size=12, color=LIGHT_GRAY)

_add_image_safe(slide, "dashboard_leaderboard", 6.0, 1.7, 6.8, 5.3)

_add_page_number(slide, 17)

# ═══════════════════════  SLIDE 18 – DASHBOARD  ══════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_slide_bg(slide, DARK_BG)
_add_title_bar(slide, "Streamlit Dashboard & SQL Chat", "Section 06")
_add_section_number_shape(slide, 0.7, 1.55, "06")

_add_textbox(slide, CONTENT_LM, 1.6, 11.5, 0.8,
             "The dashboard is a renderer, not an analyser. Every figure on every page is produced by a SQL "
             "query block loaded verbatim from sql/analysis/*.sql — the dashboard and graded deliverables "
             "cannot drift apart. pandas appears only to wrap already-aggregated rows for charting.",
             font_size=13, color=LIGHT_GRAY)

_add_image_safe(slide, "dashboard_overview", 0.7, 2.55, 6.0, 2.3)
_add_image_safe(slide, "dashboard_dataquality", 7.0, 2.55, 6.0, 2.3)

# SQL Chat
_add_textbox(slide, 0.7, 5.0, 11.5, 0.3, "SQL Chat — Text-to-SQL (Sprint 4 preview, fully local)",
             font_size=18, color=ACCENT_TEAL, bold=True)
_add_textbox(slide, CONTENT_LM, 5.35, 11.5, 1.5,
             "Type a question in English → local HuggingFace model translates to SQL → SQL runs against "
             "read-only database → result tabled and auto-charted. Safety: (1) read-only connection — "
             "writes raise at the engine, (2) whole-statement guardrail blocks destructive keywords & "
             "stacked statements, (3) execution caps: 10 s timeout + 5,000 row limit. "
             "Optional install (~2 GB), degrades gracefully to install prompt if absent.",
             font_size=12, color=LIGHT_GRAY)

_add_image_safe(slide, "sql_chat", 8.5, 5.0, 4.3, 2.1)
_add_page_number(slide, 18)

# ═══════════════════════  SLIDE 19 – DATABASE THEORY  ═════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_slide_bg(slide, DARK_BG)
_add_title_bar(slide, "Database Theory — Interview Study Guide Highlights", "Section 07")
_add_section_number_shape(slide, 0.7, 1.55, "07")

_add_textbox(slide, CONTENT_LM, 1.6, 11.8, 0.8,
             "SQL&DB_theory.md covers 6 major topic areas, each with definitions, trade-offs and worked "
             "examples measured against this project's database. Here are the key highlights:",
             font_size=13, color=LIGHT_GRAY)

theory_topics = [
    ("1. Database Paradigms", [
        "SQL vs NoSQL: schema-first, joins, ACID vs schema-flexible, horizontal scaling, eventual consistency",
        "SQLite is embedded, not server-based — zero config, single file, one writer at a time",
        "50 concurrent scrapers → SQLITE_BUSY on 84-98% of writes → migrate to PostgreSQL (MVCC)",
    ]),
    ("2. Relational Schema & Modelling", [
        "3NF: every non-key fact depends on the key, the whole key, and nothing but the key",
        "This database is 3NF with one deliberate denormalisation (departure_hour — for SARGability)",
        "PK vs FK vs UNIQUE: PK may or may NOT create a separate B-tree in SQLite (depends on type)",
    ]),
    ("3. Analytical Modelling", [
        "stop_time is the fact table (2.17M rows, event grain). All others are dimensions",
        "This schema is snowflaked: stop_time → platform → station (2-hop traversal)",
        "The grain of a fact table is the first thing you state — and the first thing that gets wrong answers",
    ]),
    ("4. ACID & CAP", [
        "ACID: Atomicity, Consistency, Isolation, Durability — 03_transform.sql runs as ONE transaction",
        "CAP: SQLite is CP (Consistent + Partition-tolerant). When API goes down, it stays consistent",
        "The 'C' in ACID ≠ the 'C' in CAP — different guarantees, shared letter by coincidence",
    ]),
    ("5. Views / Window / Subquery", [
        "Views define a word once (v_departure = 'boardable call with a time'). Without it, 5 analysts give 5 answers",
        "Window functions: aggregate WITHOUT collapsing rows. ROW_NUMBER, RANK, SUM OVER used throughout",
        "Correlated subquery runs once per outer row — usually the performance problem",
    ]),
    ("6. Index & SARGability", [
        "SARGable = Search ARGument able = the engine can use an index. Wrap a column in a function → index unusable",
        "WHERE strftime('%H', departure_time) = '17' is ~100× slower than WHERE departure_hour = 17",
        "Index Scan vs Index Seek: SEARCH cost grows with answer size, SCAN cost grows with TABLE size",
    ]),
]

for i, (title, items) in enumerate(theory_topics):
    col = i % 3
    row = i // 3
    left = 0.45 + col * 4.2
    top = 2.5 + row * 2.3
    _add_textbox(slide, left, top, 3.8, 0.3, title, font_size=12, color=ACCENT_BLUE, bold=True)
    body = "\n".join(f"• {it}" for it in items)
    _add_textbox(slide, left, top + 0.3, 3.8, 1.8, body, font_size=9, color=LIGHT_GRAY)

_add_page_number(slide, 19)

# ═══════════════════════  SLIDE 20 – TIMELINE  ═══════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_slide_bg(slide, DARK_BG)
_add_title_bar(slide, "Project Timeline", "Section 08")
_add_section_number_shape(slide, 0.7, 1.55, "08")

timeline = [
    ("Day 1", "Portal registration + API subscription. Feed reconnaissance.\n"
              "Discovered the three quirks that shaped everything:\n"
              "empty calendar.txt flags, empty wheelchair_accessible, 577K pass-through calls."),
    ("Day 2", "Schema design and the ELT pipeline — staging, 9 DQ rules,\n"
              "quarantine system, indexes, views. First full build."),
    ("Day 3", "The five analytical questions, plus the annualisation insight\n"
              "that changed Q1 and Q3 answers. Nice-to-haves:\n"
              "real-time poller, network leaderboard, index benchmarks."),
    ("Day 4", "Verification harness (21 assertions), 149-test suite,\n"
              "Streamlit dashboard, documentation, SQL Chat preview,\n"
              "study guide (SQL&DB_theory.md)."),
]

for i, (day, desc) in enumerate(timeline):
    top = 1.7 + i * 1.35
    # Circle marker
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(1.0), Inches(top + 0.05),
        Inches(0.4), Inches(0.4)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_BLUE
    shape.line.fill.background()
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = str(i+1)
    p.font.size = Pt(14)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER

    _add_textbox(slide, 1.6, top, 1.2, 0.35, day, font_size=16, color=ACCENT_BLUE, bold=True)
    _add_textbox(slide, 1.6, top + 0.32, 10.0, 0.85, desc, font_size=12, color=LIGHT_GRAY)
    # timeline connector
    if i < 3:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(1.18), Inches(top + 0.55),
            Inches(0.04), Inches(0.7)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = ACCENT_BLUE
        shape.line.fill.background()

_add_page_number(slide, 20)

# ═══════════════════════  SLIDE 21 – EVALUATION  ═════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_slide_bg(slide, DARK_BG)
_add_title_bar(slide, "Evaluation Criteria & Deliverables", "Section 08")
_add_section_number_shape(slide, 0.7, 1.55, "08")

# Must-have criteria
_add_textbox(slide, CONTENT_LM, 1.7, 5.5, 0.35, "Must-Have Criteria",
             font_size=18, color=ACCENT_BLUE, bold=True)
must_have = [
    "SQLite database with normalised schema and Foreign Keys ✓",
    "Schema diagram (ERD + drawDB) ✓",
    "Clean data in every table (9 DQ rules, 12 quarantined of 7M rows) ✓",
    "All 5 core analytical questions accurately answered (SQL only) ✓",
    "Table definitions + queries in dedicated .sql files ✓",
    "Visualisation components match data queries (Streamlit dashboard) ✓",
]
for i, item in enumerate(must_have):
    _add_textbox(slide, CONTENT_LM, 2.15 + i * 0.35, 5.5, 0.3, f"• {item}", font_size=11, color=LIGHT_GRAY)

# Nice-to-have
_add_textbox(slide, CONTENT_LM, 4.4, 5.5, 0.35, "Nice-to-Have (all completed)",
             font_size=18, color=ACCENT_TEAL, bold=True)
nice_have = [
    "Live Stream Integration (GTFS-RT poller, cron/launchd) ✓",
    "Network Leaderboard (5 hubs compared, structural + punctuality) ✓",
    "Index Optimisation (EXPLAIN QUERY PLAN, benchmarks, make benchmark) ✓",
    "Team Study Guide (SQL&DB_theory.md — 6 major topic areas, 1,000+ lines) ✓",
    "SQL Chat — Text-to-SQL GenAI preview (local model, safety guardrails) ✓",
]
for i, item in enumerate(nice_have):
    _add_textbox(slide, CONTENT_LM, 4.85 + i * 0.35, 5.5, 0.3, f"• {item}", font_size=11, color=LIGHT_GRAY)

# Deliverables checklist
_add_textbox(slide, 6.8, 1.7, 5.8, 0.35, "Deliverables",
             font_size=18, color=ACCENT_AMBER, bold=True)
deliverables = [
    "GitHub repository with source code",
    "README.md — description, ERD, visuals, contributors, timeline, personal statement",
    "team feedback: 5-minute overview + approach, DB design, answers",
    "Q&A: random question from SQL&DB_theory.md study guide",
    "PowerPoint presentation (this file!)",
]
for i, item in enumerate(deliverables):
    _add_textbox(slide, 6.8, 2.15 + i * 0.45, 5.9, 0.38, f"• {item}", font_size=12, color=LIGHT_GRAY)

# Stats block
_add_textbox(slide, 6.8, 4.5, 5.9, 0.35, "By the Numbers",
             font_size=18, color=ACCENT_BLUE, bold=True)
stats = [
    "7,057,090 rows staged → 12 quarantined",
    "149 tests · 47 labelled analysis queries",
    "34 enforced foreign keys · 21 verify assertions",
    "~3 min from ZIP to verified, indexed database",
    "~1,750 lines of Python · ~3,000+ lines of SQL",
    "PEP 517 package, editable install, Makefile",
]
for i, stat in enumerate(stats):
    _add_textbox(slide, 6.8, 4.95 + i * 0.35, 5.9, 0.3, f"• {stat}", font_size=11, color=LIGHT_GRAY)

_add_page_number(slide, 21)

# ═══════════════════════  SLIDE 22 – CONCLUSION ══════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_slide_bg(slide, DARK_BG)

# accent bar at top
shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.08)
)
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT_BLUE
shape.line.fill.background()

_add_textbox(slide, 0.7, 0.8, 11.5, 1.0, "Thank You",
             font_size=48, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
_add_accent_line(slide, 5.5, 1.8, 2.3)

_add_textbox(slide, 0.7, 2.3, 11.5, 0.6,
             "RailPulse — Belgian Transit SQL Analysis",
             font_size=24, color=ACCENT_BLUE, bold=False, alignment=PP_ALIGN.CENTER)

_add_textbox(slide, 0.7, 3.1, 11.5, 2.0,
             "The interesting part of this project was not the SQL — it was discovering that "
             "the two most natural queries give the wrong answer.\n\n"
             "The query is the easy part. Knowing what the data actually is — and being willing "
             "to publish the number you did NOT pick, next to the one you did — is the job.",
             font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

_add_textbox(slide, 0.7, 5.2, 11.5, 0.8,
             "Stéphane van der Aa  ·  BeCode AI & Data Science  ·  July 2026\n"
             "Data: NMBS-SNCB – Open Data – 2026-07-20  ·  Licensed CC BY 4.0",
             font_size=13, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

_add_page_number(slide, 22)

# ══════════════════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════════════════
prs.save(OUTPUT_PATH)
print(f"✓ Presentation saved to {OUTPUT_PATH}")
print(f"  Slides: {len(prs.slides)}")
